from flask import Flask, render_template, redirect, url_for, request, send_file, url_for, jsonify, session, flash
import imaplib
import email
from email.header import decode_header
from bs4 import BeautifulSoup
import re
import dateparser
import os
import pdfkit
from werkzeug.utils import secure_filename
import base64
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from openai import OpenAI
from dotenv import load_dotenv
import os

load_dotenv()  # Load environment variables from .env file

app = Flask(__name__)
app.secret_key = os.getenv('FLASK_SECRET_KEY', 'fallback_secret_key')

EMAIL = os.getenv('SENDER_EMAIL')
PASSWORD = os.getenv('SENDER_APP_PASSWORD')

import random

# Store OTPs temporarily in memory (for demo purposes)
otp_store = {}

def generate_otp():
    return str(random.randint(100000, 999999))

# Configure OpenAI client
openai_api_key = os.getenv('OPENAI_API_KEY')
if not openai_api_key:
    raise ValueError("OpenAI API key not set in environment variables")

client = OpenAI(
    base_url="https://models.inference.ai.azure.com",
    api_key=openai_api_key
)

import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
import os

def send_otp_email(to_email, otp):
    sender_email = os.getenv("SENDER_EMAIL")
    sender_password = os.getenv("SENDER_APP_PASSWORD")
    if not sender_email or not sender_password:
        raise ValueError("Sender email credentials not set in environment variables")

    message = MIMEMultipart("alternative")
    message["Subject"] = "Your OTP Verification Code"
    message["From"] = sender_email
    message["To"] = to_email

    text = f"Your OTP code is: {otp}"
    part = MIMEText(text, "plain")
    message.attach(part)

    with smtplib.SMTP_SSL("smtp.gmail.com", 465) as server:
        server.login(sender_email, sender_password)
        server.sendmail(sender_email, to_email, message.as_string())

# Configure upload folder
UPLOAD_FOLDER = 'static/uploads'
ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif'}
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Create upload folder if it doesn't exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Path to wkhtmltopdf executable
config = pdfkit.configuration(
    wkhtmltopdf=r"C:\Program Files\wkhtmltopdf\bin\wkhtmltopdf.exe",
    # wkhtmltopdf=r"/usr/local/bin/wkhtmltopdf",

)

def categorize(subject, body):
    combined = (subject + " " + body).lower()
    if "internship" in combined:
        return "Internship"
    elif "full time" in combined or "full-time" in combined or "job" in combined:
        return "Full-time Job"
    elif "hackathon" in combined:
        return "Hackathon"
    elif any(k in combined for k in ["apply", "register", "opportunity", "opening", "recruitment", "role"]):
        return "Other Opportunities"
    else:
        return "Ignore"

def clean_email_body(raw_html):
    soup = BeautifulSoup(raw_html, 'html.parser')
    for tag in soup(["style", "script", "head", "meta", "title"]):
        tag.decompose()
    body = soup.body or soup
    text = body.get_text(separator="\n").strip()

    keywords = ['deadline', 'apply', 'CTC', 'salary', 'last date', 'job', 'internship']
    for kw in keywords:
        regex = re.compile(rf'\b({kw})\b', re.IGNORECASE)
        text = regex.sub(
            r'<span style="background: #fff3cd; color: #a76f00; font-weight: bold;">\1</span>', text)
    text = re.sub(r'(https?://[^\s]+)',
                  r'<a href="\1" target="_blank" style="color:#4b0082; text-decoration:underline;">\1</a>', text)
    return text

def fetch_emails():
    results = []
    imap = imaplib.IMAP4_SSL("imap.gmail.com")
    imap.login(EMAIL, PASSWORD)
    imap.select("inbox")

    status, messages = imap.search(None, "ALL")
    email_ids = messages[0].split()

    for eid in email_ids[-50:]:
        _, msg_data = imap.fetch(eid, "(RFC822)")
        raw_email = msg_data[0][1]
        msg = email.message_from_bytes(raw_email)

        subject, encoding = decode_header(msg["Subject"])[0]
        if isinstance(subject, bytes):
            subject = subject.decode(encoding or "utf-8", errors="ignore")
        from_ = msg.get("From")
        received = msg.get("Date", "Unknown")
        body = ""
        plain_text_body = ""
        app_link = "Not provided"
        deadline = "Not found"
        eligibility = "Not mentioned"

        if msg.is_multipart():
            for part in msg.walk():
                ctype = part.get_content_type()
                if ctype == "text/html":
                    raw_html = part.get_payload(decode=True).decode(errors="ignore")
                    soup = BeautifulSoup(raw_html, "html.parser")
                    plain_text_body = soup.get_text(separator="\n").strip()
                    body = clean_email_body(raw_html)
                    break
                elif ctype == "text/plain" and not plain_text_body:
                    plain_text_body = part.get_payload(decode=True).decode(errors="ignore")
                    body = clean_email_body(plain_text_body)
        else:
            raw_text = msg.get_payload(decode=True).decode(errors="ignore")
            plain_text_body = BeautifulSoup(raw_text, "html.parser").get_text(separator="\n").strip()
            body = clean_email_body(raw_text)

        # Extract apply link
        link_match = re.search(r'https?://[a-zA-Z0-9./?=_%-]+', body)
        if link_match:
            extracted_link = link_match.group(0)
            if "http" in extracted_link and "." in extracted_link:
                app_link = extracted_link

        # Deadline Extraction: keyword-based first
        found = False
        for line in plain_text_body.splitlines():
            for kw in ["deadline", "apply before", "last date", "register by", "before"]:
                if kw in line.lower():
                    parsed = dateparser.parse(line, settings={'PREFER_DATES_FROM': 'future'})
                    if parsed:
                        deadline = parsed.strftime("%d/%m/%Y")
                        found = True
                        break
            if found:
                break

        # If no keyword-based date, try any date
        if deadline == "Not found":
            for line in plain_text_body.splitlines():
                parsed = dateparser.parse(line, settings={'PREFER_DATES_FROM': 'future'})
                if parsed:
                    deadline = parsed.strftime("%d/%m/%Y")
                    break

        # Eligibility
        elig_match = re.search(r'(eligibility[^:\n]*[:\n])(.+)', body, re.IGNORECASE)
        if elig_match:
            eligibility = elig_match.group(2).strip()

        category = categorize(subject, body)

        if category != "Ignore":
            results.append({
                "title": subject.strip(),
                "from": from_,
                "received": received,
                "description": body,
                "deadline": deadline,
                "eligibility": eligibility,
                "application_link": app_link,
                "category": category
            })
            print(f"[DEBUG] {subject} → Deadline: {deadline}")

    imap.logout()
    return results

email_cache = []

@app.route('/')
def home():
    if 'email' not in session or 'password' not in session:
        return redirect(url_for('register'))
    return render_template('landing.html')

@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        email_input = request.form.get('email')
        app_password = request.form.get('app_password')
        if not email_input or not app_password:
            flash("Email and App Password are required.", "error")
            return render_template('registration.html')
        otp = generate_otp()
        otp_store[email_input] = otp
        try:
            send_otp_email(email_input, otp)
        except Exception as e:
            flash(f"Failed to send OTP email: {str(e)}", "error")
            return render_template('registration.html')
        session['pending_email'] = email_input
        session['pending_password'] = app_password
        return redirect(url_for('verify_otp'))
    return render_template('registration.html')

@app.route('/verify_otp', methods=['GET', 'POST'])
def verify_otp():
    if request.method == 'POST':
        user_otp = request.form.get('otp')
        pending_email = session.get('pending_email')
        if not pending_email or not user_otp:
            flash("Invalid session or OTP.", "error")
            return redirect(url_for('register'))
        real_otp = otp_store.get(pending_email)
        if real_otp == user_otp:
            session['email'] = pending_email
            session['password'] = session.get('pending_password')
            otp_store.pop(pending_email, None)
            session.pop('pending_email', None)
            session.pop('pending_password', None)
            return redirect(url_for('home'))
        else:
            flash("Incorrect OTP. Please try again.", "error")
    return render_template('otp_verification.html')

@app.route('/analyze')
def analyze():
    global email_cache
    email_cache = fetch_emails()
    return redirect(url_for('show_categories'))

@app.route('/categories')
def show_categories():
    predefined = ["Internship", "Full-time Job", "Hackathon", "Other Opportunities"]
    category_counts = {cat: 0 for cat in predefined}

    for email in email_cache:
        cat = email["category"]
        if cat in category_counts:
            category_counts[cat] += 1
        else:
            category_counts["Other Opportunities"] += 1

    total = len(email_cache)
    return render_template("categories.html", category_counts=category_counts, total=total)

@app.route('/opportunities/<category>')
def show_opportunities(category):
    filtered = [e for e in email_cache if e["category"] == category]
    return render_template("opportunities.html", opportunities=filtered, category=category)

# Resume functionality integration

from flask import request, send_file, jsonify

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in {'png', 'jpg', 'jpeg', 'gif'}

def generate_role_suggestions(role):
    try:
        response = client.chat.completions.create(
            messages=[
                {
                    "role": "system",
                    "content": "You are a professional resume writer. Generate 3 different professional summaries for the given job role. Each summary should be 2-3 sentences long and highlight relevant skills and experience."
                },
                {
                    "role": "user",
                    "content": f"Generate professional summary for a {role} position."
                }
            ],
            model="gpt-4o-mini",
            temperature=0.7,
            max_tokens=300,
            top_p=1
        )
        suggestions = [choice.message.content.strip() for choice in response.choices]
        return suggestions
    except Exception as e:
        print(f"Error generating role suggestions: {str(e)}")
        return []

def generate_skill_suggestions(role):
    try:
        response = client.chat.completions.create(
            messages=[
                {
                    "role": "system",
                    "content": "You are a professional resume writer. Generate 5-7 relevant technical and soft skills for the given job role. Format each skill as a single word or short phrase without numbering."
                },
                {
                    "role": "user",
                    "content": f"Generate relevant skills for a {role} position without numbering."
                }
            ],
            model="gpt-4o-mini",
            temperature=0.7,
            max_tokens=200,
            top_p=1
        )
        suggestions = [skill.strip() for skill in response.choices[0].message.content.split('\n') if skill.strip()]
        return suggestions
    except Exception as e:
        print(f"Error generating skill suggestions: {str(e)}")
        return []

def create_presentation(data):
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = data['name']
    subtitle.text = data['title']
    
    # About slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "About Me"
    tf = body_shape.text_frame
    tf.text = data['about']
    
    # Education slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Education"
    tf = body_shape.text_frame
    
    if data.get('edu1'):
        p = tf.add_paragraph()
        p.text = data['edu1']
    if data.get('edu2'):
        p = tf.add_paragraph()
        p.text = data['edu2']
    if data.get('edu3'):
        p = tf.add_paragraph()
        p.text = data['edu3']
    
    # Experience slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Experience"
    tf = body_shape.text_frame
    
    if data.get('exp1'):
        p = tf.add_paragraph()
        p.text = data['exp1']
    if data.get('exp2'):
        p = tf.add_paragraph()
        p.text = data['exp2']
    
    # Skills slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Skills"
    tf = body_shape.text_frame
    
    if data.get('skills'):
        for skill in data['skills'].split(','):
            p = tf.add_paragraph()
            p.text = skill.strip()
    
    # References slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "References"
    tf = body_shape.text_frame
    
    if data.get('ref1'):
        p = tf.add_paragraph()
        p.text = data['ref1']
    if data.get('ref2'):
        p = tf.add_paragraph()
        p.text = data['ref2']
    
    # Save the presentation
    presentation_path = "resume_presentation.pptx"
    prs.save(presentation_path)
    return presentation_path

@app.route('/resume', methods=['GET', 'POST'])
def resume_form():
    if request.method == 'POST':
        data = request.form.to_dict()
        
        # Handle image upload
        if 'profile_image' in request.files:
            file = request.files['profile_image']
            if file and allowed_file(file.filename):
                filename = secure_filename(file.filename)
                filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                file.save(filepath)
                data['profile_image'] = url_for('static', filename=f'uploads/{filename}')
        
        # Choose template based on format
        template = 'resume/resume_professional.html' if data.get('format') == 'professional' else 'resume/resume.html'
        
        return render_template(template, data=data)
    return render_template('resume/form.html')

@app.route('/resume/get_role_suggestions', methods=['POST'])
def resume_get_role_suggestions():
    data = request.get_json()
    role = data.get('role', '')
    suggestions = generate_role_suggestions(role)
    return jsonify({'suggestions': suggestions})

@app.route('/resume/get_skill_suggestions', methods=['POST'])
def resume_get_skill_suggestions():
    data = request.get_json()
    role = data.get('role', '')
    suggestions = generate_skill_suggestions(role)
    return jsonify({'suggestions': suggestions})

@app.route('/resume/download', methods=['POST'])
def resume_download():
    data = request.form.to_dict()
    
    # Handle image in the data
    if 'profile_image' in data:
        # Convert image to base64 for PDF embedding
        image_path = os.path.join(app.root_path, 'static', data['profile_image'].lstrip('/'))
        if os.path.exists(image_path):
            with open(image_path, 'rb') as img_file:
                data['profile_image'] = f"data:image/jpeg;base64,{base64.b64encode(img_file.read()).decode()}"
    
    # Choose template based on format
    template = 'resume/resume_professional.html' if data.get('format') == 'professional' else 'resume/resume.html'
    html = render_template(template, data=data)

    # Enable local file access and set PDF options
    options = {
        'enable-local-file-access': '',
        'page-size': 'A4',
        'margin-top': '0.75in',
        'margin-right': '0.75in',
        'margin-bottom': '0.75in',
        'margin-left': '0.75in',
        'encoding': 'UTF-8',
        'no-outline': None,
        'quiet': ''
    }

    # Generate the PDF
    pdfkit.from_string(html, 'resume.pdf', configuration=config, options=options)
    
    # If presentation format is requested
    if data.get('output_presentation'):
        presentation_path = create_presentation(data)
        return jsonify({
            'pdf_url': url_for('resume_download_pdf'),
            'presentation_url': url_for('resume_download_presentation')
        })
    
    return send_file('resume.pdf', as_attachment=True)

@app.route('/resume/download_pdf')
def resume_download_pdf():
    return send_file('resume.pdf', as_attachment=True)

@app.route('/resume/download_presentation')
def resume_download_presentation():
    return send_file('resume_presentation.pptx', as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)
